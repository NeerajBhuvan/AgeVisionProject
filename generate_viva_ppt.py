"""Generate the AgeVision Viva-Voce deck — 20 slides, Anna University (Chennai) format.

Content is aligned strictly to THESIS_REPORT_NEERAJ_BHUVAN: 3 progression engines,
13 modules, 21 REST endpoints, 6 MongoDB collections, 50 test cases, MAE 3.65 yr.
Plain, easy English. Font: Times New Roman. No em/en dashes in slide text.

Outputs:
    AgeVision_Viva_AnnaUniv.pptx
    AgeVision_Viva_AnnaUniv.pdf   (via PowerPoint COM automation)

Run: python generate_viva_ppt.py
"""
from __future__ import annotations

import os
import subprocess
import sys

try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree

try:
    from PIL import Image
    _HAS_PIL = True
except ImportError:
    _HAS_PIL = False

# ── paths ──────────────────────────────────────────────────────────────────
ROOT     = os.path.dirname(os.path.abspath(__file__))
SS_DIR   = os.path.join(ROOT, "screenshots")
APP_LOGO = os.path.join(ROOT, "app_menu_logo.png")
OUT      = os.path.join(ROOT, "AgeVision_Viva_AnnaUniv.pptx")
PDF      = os.path.join(ROOT, "AgeVision_Viva_AnnaUniv.pdf")

DATE_LABEL = "June 2026"
TOTAL = 20


def ss(name: str) -> str:
    return os.path.join(SS_DIR, name)


# ── slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── colour palette — white / light theme ─────────────────────────────────────
BG      = RGBColor(0xFF, 0xFF, 0xFF)
CARD    = RGBColor(0xF4, 0xF1, 0xFC)
CARD2   = RGBColor(0xF9, 0xF7, 0xFF)
BORDER  = RGBColor(0xD4, 0xC8, 0xF5)
PURPLE  = RGBColor(0x7F, 0x5A, 0xF0)   # --accent1 (single accent)
GREEN   = RGBColor(0x7F, 0x5A, 0xF0)   # retired: aliased to purple (monochrome)
ORANGE  = RGBColor(0x7F, 0x5A, 0xF0)   # retired: aliased to purple
CHIP    = RGBColor(0xED, 0xE8, 0xFC)   # soft lavender fill (low-ink)
WHITE   = RGBColor(0xFF, 0xFF, 0xFE)
MUTED   = RGBColor(0x3D, 0x35, 0x53)   # --text-secondary
TEXT    = RGBColor(0x1A, 0x16, 0x25)   # --text
HDR_END = RGBColor(0x6C, 0x3B, 0xD5)   # softened deep purple (header gradient end)
NAVY    = RGBColor(0x1A, 0x3A, 0x5C)
GOLD    = RGBColor(0xC8, 0x92, 0x2A)

FONT = "Times New Roman"


# ═════════════════════════════════════════════════════════════════════════ #
#                              HELPER FUNCTIONS                              #
# ═════════════════════════════════════════════════════════════════════════ #

def _hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, *, border_color=None, border_pt=0.75,
             rounded=False):
    auto = 5 if rounded else 1
    shape = slide.shapes.add_shape(auto, l, t, w, h)
    shape.shadow.inherit = False
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_gradient(slide, l, t, w, h, c_start=PURPLE, c_end=HDR_END, angle=0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.shadow.inherit = False
    shape.line.fill.background()
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    grad = (
        f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{_hex(c_start)}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{_hex(c_end)}"/></a:gs>'
        f'</a:gsLst><a:lin ang="{angle}" scaled="0"/></a:gradFill>'
    )
    spPr.insert(0, etree.fromstring(grad))
    return shape


def add_text(slide, l, t, w, h, text, *, size=14, bold=False, italic=False,
             color=TEXT, align=PP_ALIGN.LEFT, font=FONT, wrap=True,
             anchor=None, line_spacing=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tf


def add_runs(slide, l, t, w, h, runs, *, size=22, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, font=FONT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for txt, col, bold in runs:
        r = p.add_run()
        r.text = txt
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col
    return tf


def add_bullets(slide, l, t, w, h, items, *, size=13, color=TEXT, font=FONT,
                bullet="•  ", space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            lead, rest = item
            r0 = p.add_run(); r0.text = bullet + lead
            r0.font.name = font; r0.font.size = Pt(size); r0.font.bold = True
            r0.font.color.rgb = color
            r1 = p.add_run(); r1.text = rest
            r1.font.name = font; r1.font.size = Pt(size); r1.font.bold = False
            r1.font.color.rgb = color
        else:
            r = p.add_run(); r.text = bullet + item
            r.font.name = font; r.font.size = Pt(size); r.font.bold = False
            r.font.color.rgb = color
    return tf


def add_divider(slide, l, t, w, color=PURPLE, pt=2):
    return add_rect(slide, l, t, w, Pt(pt), color)


def _img_size(path, max_w_emu, max_h_emu):
    if _HAS_PIL and os.path.exists(path):
        iw, ih = Image.open(path).size
        ar = iw / ih
    else:
        ar = 1.85
    box_ar = max_w_emu / max_h_emu
    if ar >= box_ar:
        w = max_w_emu
        h = int(max_w_emu / ar)
    else:
        h = max_h_emu
        w = int(max_h_emu * ar)
    return int(w), int(h)


def add_image_fit(slide, path, l, t, max_w, max_h, *, border=None, border_pt=1.25,
                  halign="center", valign="middle"):
    if not os.path.exists(path):
        add_rect(slide, l, t, max_w, max_h, CARD, border_color=BORDER)
        add_text(slide, l, t, max_w, max_h, os.path.basename(path),
                 size=9, color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return None
    w, h = _img_size(path, int(max_w), int(max_h))
    if halign == "center":
        ox = l + (int(max_w) - w) // 2
    elif halign == "right":
        ox = l + (int(max_w) - w)
    else:
        ox = l
    if valign == "middle":
        oy = t + (int(max_h) - h) // 2
    elif valign == "bottom":
        oy = t + (int(max_h) - h)
    else:
        oy = t
    pic = slide.shapes.add_picture(path, Emu(int(ox)), Emu(int(oy)),
                                   width=Emu(w), height=Emu(h))
    if border:
        frame = add_rect(slide, Emu(int(ox)), Emu(int(oy)), Emu(w), Emu(h), None,
                         border_color=border, border_pt=border_pt)
        frame.shadow.inherit = False
    return pic


def header(slide, title, num, accent=PURPLE):
    accent = PURPLE  # unified header colour across all slides
    add_gradient(slide, 0, 0, W, Inches(0.82), accent, HDR_END, angle=0)
    add_text(slide, Inches(0.45), 0, Inches(11.4), Inches(0.82), title,
             size=25, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.7), 0, Inches(1.4), Inches(0.82),
             f"{num:02d} / {TOTAL}", size=13, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 0, Inches(7.34), W, Inches(0.16), accent)
    add_text(slide, Inches(0.45), Inches(7.16), Inches(12.4), Inches(0.18),
             f"AgeVision  •  MCA (Final Semester) Project  •  Anna University, CDE Chennai  •  {DATE_LABEL}",
             size=9, color=MUTED)


def section_label(slide, l, t, w, text, color):
    add_text(slide, l, t, w, Inches(0.34), text, size=15, bold=True, color=color)
    add_divider(slide, l, t + Inches(0.4), w, color, 1.5)


def card(slide, l, t, w, h, accent=BORDER, fill=CARD, border_pt=0.75):
    return add_rect(slide, l, t, w, h, fill, border_color=accent, border_pt=border_pt)


def stat_chips(slide, y, chips, *, x0=Inches(0.4), total_w=Inches(12.55),
               h=Inches(0.92), gap=Inches(0.18)):
    """chips: list of (big, small, color)."""
    n = len(chips)
    cw = int((int(total_w) - int(gap) * (n - 1)) / n)
    for i, (big, small, col) in enumerate(chips):
        x = int(x0) + i * (cw + int(gap))
        add_rect(slide, Emu(x), y, Emu(cw), h, CHIP, border_color=PURPLE,
                 border_pt=1.0, rounded=True)
        add_text(slide, Emu(x), y + Inches(0.12), Emu(cw), Inches(0.42), big,
                 size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
        add_text(slide, Emu(x), y + Inches(0.55), Emu(cw), Inches(0.3), small,
                 size=10.5, color=MUTED, align=PP_ALIGN.CENTER)


def demo_grid(slide, items, *, top=Inches(1.05)):
    """items: list of (path, caption, color), up to 4, laid out 2x2."""
    cell_w, cell_h = Inches(6.18), Inches(2.4)
    gx, gy, cap_h = Inches(0.19), Inches(0.18), Inches(0.3)
    x0 = Inches(0.4)
    for i, (path, cap, col) in enumerate(items):
        c, r = i % 2, i // 2
        x = x0 + c * (cell_w + gx)
        y = top + r * (cell_h + cap_h + gy)
        add_image_fit(slide, path, x, y, cell_w, cell_h, border=col, valign="middle")
        add_text(slide, x, y + cell_h + Inches(0.01), cell_w, cap_h, cap,
                 size=11, bold=True, color=col, align=PP_ALIGN.CENTER)


# ── table helpers ────────────────────────────────────────────────────────────

def _cell_border(cell, color=BORDER):
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tcPr = cell._tc.get_or_add_tcPr()
    for side in ("lnL", "lnR", "lnT", "lnB"):
        ln = etree.SubElement(tcPr, f"{{{NS}}}{side}")
        ln.set("w", "9525")
        sf = etree.SubElement(ln, f"{{{NS}}}solidFill")
        etree.SubElement(sf, f"{{{NS}}}srgbClr").set("val", _hex(color))


def fill_table(tbl, rows, *, header_fill=PURPLE, odd=RGBColor(0xED, 0xE8, 0xFC),
               even=RGBColor(0xF8, 0xF6, 0xFF), hdr_pt=12, body_pt=11,
               last_fill=None, align_first_left=True):
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
    nrows = len(rows)
    for r in range(nrows):
        is_last = last_fill is not None and r == nrows - 1
        for c in range(len(rows[0])):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = header_fill
            elif is_last:
                cell.fill.fore_color.rgb = last_fill
            else:
                cell.fill.fore_color.rgb = odd if r % 2 == 1 else even
            for p in cell.text_frame.paragraphs:
                p.alignment = (PP_ALIGN.LEFT if (c == 0 and align_first_left and r != 0)
                               else PP_ALIGN.CENTER)
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(hdr_pt if r == 0 else body_pt)
                    run.font.bold = (r == 0 or is_last)
                    run.font.color.rgb = (WHITE if (r == 0 or is_last) else TEXT)
            _cell_border(cell)


# ═════════════════════════════════════════════════════════════════════════ #
#                                  SLIDES                                    #
# ═════════════════════════════════════════════════════════════════════════ #

def s01_title(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    add_gradient(slide, 0, 0, W, Inches(0.16), PURPLE, HDR_END)
    add_gradient(slide, 0, Inches(7.34), W, Inches(0.16), HDR_END, PURPLE)

    add_text(slide, Inches(0.5), Inches(0.34), Inches(12.33), Inches(0.42),
             "ANNA UNIVERSITY, CHENNAI 600 025", size=20, bold=True,
             color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(0.78), Inches(12.33), Inches(0.32),
             "CENTRE FOR DISTANCE EDUCATION", size=13, bold=True,
             color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(1.10), Inches(12.33), Inches(0.3),
             "Faculty of Information and Communication Engineering", size=12,
             color=MUTED, align=PP_ALIGN.CENTER)

    if os.path.exists(APP_LOGO):
        add_image_fit(slide, APP_LOGO, Inches(4.62), Inches(1.5), Inches(1.05),
                      Inches(0.72), halign="left", valign="middle")
    add_runs(slide, Inches(5.78), Inches(1.5), Inches(4.0), Inches(0.72),
             [("AgeVision", NAVY, True), (" AI", GOLD, True)],
             size=26, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.5), Inches(2.46), Inches(12.33), Inches(0.9),
             "AI Based Age Prediction and Age Progression\nSystem Using Facial Images",
             size=31, bold=True, color=TEXT, align=PP_ALIGN.CENTER, line_spacing=1.05)

    add_text(slide, Inches(0.5), Inches(3.66), Inches(12.33), Inches(0.34),
             "A Project Report submitted in partial fulfilment for the award of the degree of",
             size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.36),
             "MASTER OF COMPUTER APPLICATIONS", size=16, bold=True,
             color=GREEN, align=PP_ALIGN.CENTER)

    add_divider(slide, Inches(4.4), Inches(4.56), Inches(4.53), PURPLE, 2.5)

    cw = Inches(5.9)
    card(slide, Inches(0.65), Inches(4.78), cw, Inches(1.7), accent=BORDER)
    add_text(slide, Inches(0.9), Inches(4.94), cw - Inches(0.5), Inches(0.3),
             "PRESENTED BY", size=12, bold=True, color=PURPLE)
    add_bullets(slide, Inches(0.9), Inches(5.3), cw - Inches(0.5), Inches(1.1), [
        ("Name:  ", "Neeraj Bhuvan M"),
        ("Roll No:  ", "2435MCA0011"),
        ("Register No:  ", "67224100038"),
    ], size=14, bullet="", space_after=4)

    card(slide, Inches(6.78), Inches(4.78), cw, Inches(1.7), accent=BORDER)
    add_text(slide, Inches(7.03), Inches(4.94), cw - Inches(0.5), Inches(0.3),
             "PROJECT GUIDE", size=12, bold=True, color=GREEN)
    add_bullets(slide, Inches(7.03), Inches(5.3), cw - Inches(0.5), Inches(1.1), [
        ("Dr. P. Geetha", ""),
        ("Associate Professor", ""),
        ("Dept. of Information Science & Technology, CEG", ""),
    ], size=13.5, bullet="", space_after=3)

    add_text(slide, Inches(0.5), Inches(6.64), Inches(12.33), Inches(0.34),
             f"Viva-Voce Examination  •  {DATE_LABEL}", size=13, bold=True,
             color=ORANGE, align=PP_ALIGN.CENTER)


def s02_abstract(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Abstract", 2, PURPLE)

    add_rect(slide, Inches(0.4), Inches(1.0), Inches(12.55), Inches(1.25), CARD2,
             border_color=BORDER, border_pt=1.0)
    add_text(slide, Inches(0.62), Inches(1.12), Inches(12.1), Inches(1.05),
             "AgeVision is a browser-based, AI-driven platform that brings automated age estimation and "
             "identity-preserving facial age progression together in one secure interface. The client uses "
             "Angular 19; the server uses Django 5.2 with DRF 3.16, protected by JWT, Bcrypt and Fernet. "
             "A hybrid store keeps AI results in MongoDB (six collections) and logins in SQLite.",
             size=13, color=TEXT, line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)

    add_bullets(slide, Inches(0.4), Inches(2.5), Inches(6.15), Inches(2.7), [
        "Prediction uses two models together: MiVOLO v2 (about 3.65 yr error) with an InsightFace backup (about 8.5 yr).",
        "This is about 48 to 54% more accurate than How-Old.net and Face++ (6 to 8 yr).",
        "Progression uses three engines: SAM GAN, Fast-AgingGAN and FADING diffusion.",
    ], size=12.5, color=TEXT, space_after=12, line_spacing=1.08)

    add_bullets(slide, Inches(6.8), Inches(2.5), Inches(6.15), Inches(2.7), [
        "Every aged face passes a FaceNet identity check (score above 0.6), met by over 95% of samples.",
        "The platform has 13 modules, 21 REST endpoints and 7-class emotion reading.",
        "It works as a reusable base for forensic, healthcare, entertainment, security and research use.",
    ], size=12.5, color=TEXT, space_after=12, line_spacing=1.08)

    stat_chips(slide, Inches(5.95), [
        ("3.65 yr", "average age error", PURPLE),
        ("> 95%", "keep identity", GREEN),
        ("13", "modules", PURPLE),
        ("21", "REST endpoints", GREEN),
    ])


def s03_intro(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Introduction & Problem Statement", 3, PURPLE)

    section_label(slide, Inches(0.4), Inches(1.0), Inches(6.05), "Background", PURPLE)
    add_bullets(slide, Inches(0.4), Inches(1.55), Inches(6.05), Inches(3.6), [
        "More fields now use computers to read faces, such as apps, hospitals and police work.",
        "Two needs keep coming up: guess a person's age from a photo, and show how the face may look at another age while keeping the same identity.",
        "Human guesses differ from person to person, paid cloud tools cost money and raise privacy worries, and most free tools need developer skills.",
        "Doing both jobs well, for free, in one website, was not available before.",
    ], size=12.5, color=TEXT, space_after=12, line_spacing=1.06)

    section_label(slide, Inches(6.75), Inches(1.0), Inches(6.2), "Limits of existing tools", GREEN)
    add_bullets(slide, Inches(6.75), Inches(1.55), Inches(6.2), Inches(3.6), [
        ("How-Old.net and DeepFace", " only guess age; they cannot age a face."),
        ("Face++", " is accurate but paid, and runs only in the cloud."),
        ("AgingBooth", " can age a face but is mobile-only and does not check identity."),
        ("Open-source libraries", " mostly do prediction only and need developer set-up."),
    ], size=12.5, color=TEXT, space_after=12, line_spacing=1.06)

    add_rect(slide, Inches(0.4), Inches(5.35), Inches(12.55), Inches(1.65), CARD2,
             border_color=GREEN, border_pt=1.0)
    add_text(slide, Inches(0.62), Inches(5.5), Inches(12.1), Inches(0.3),
             "The gap AgeVision fills", size=13.5, bold=True, color=GREEN)
    add_text(slide, Inches(0.62), Inches(5.88), Inches(12.1), Inches(1.0),
             "AgeVision is one free, self-hosted web app that both guesses age and ages a face, keeps every "
             "photo on your own computer, and adds group, batch, emotion, live camera, history, analytics and "
             "an admin panel in a single login.",
             size=12.5, color=TEXT, line_spacing=1.1)


def s04_objectives(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Objectives & Scope", 4, GREEN)

    section_label(slide, Inches(0.4), Inches(1.0), Inches(7.55), "Objectives", PURPLE)
    objectives = [
        "Build a web app (Angular 19 + Django 5.2) that guesses and ages faces from one photo, many photos, or a live camera.",
        "Guess age using two models together: MiVOLO v2 (main, about 3.65 yr error) and InsightFace (backup, about 8.5 yr error).",
        "Offer three ways to age a face: SAM GAN (main), Fast-AgingGAN (runs on CPU), FADING diffusion (ages and de-ages; optional cloud GPU).",
        "Detect 7 emotions on every face along with the age.",
        "Keep the app safe with JWT, Bcrypt and Fernet, and show live progress for slow jobs.",
        "Store data in two databases working together: MongoDB and SQLite.",
        "Add an admin-only page with charts, user control, and live health checks.",
    ]
    y = Inches(1.55)
    for i, obj in enumerate(objectives, 1):
        add_rect(slide, Inches(0.4), y + Inches(0.02), Inches(0.34), Inches(0.34),
                 PURPLE if i % 2 else GREEN, rounded=True)
        add_text(slide, Inches(0.4), y + Inches(0.02), Inches(0.34), Inches(0.34),
                 str(i), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(0.86), y, Inches(7.1), Inches(0.74), obj, size=12.5,
                 color=TEXT, line_spacing=0.98)
        y += Inches(0.79)

    card(slide, Inches(8.2), Inches(1.0), Inches(4.75), Inches(6.05), accent=GREEN, border_pt=1.5)
    add_text(slide, Inches(8.42), Inches(1.18), Inches(4.3), Inches(0.34), "Scope",
             size=15, bold=True, color=GREEN)
    add_divider(slide, Inches(8.42), Inches(1.58), Inches(4.3), GREEN, 1.5)
    add_bullets(slide, Inches(8.42), Inches(1.74), Inches(4.3), Inches(5.1), [
        "Sign up, log in, reset password, edit profile.",
        "Guess age for one face, a group photo, or many photos at once.",
        "Detect emotion with every age guess.",
        "Age a face while keeping the same identity.",
        "Personal history, charts dashboard and settings.",
        "Admin page for users, charts and system health.",
        "API can be used by mobile apps or other systems.",
    ], size=12.5, color=TEXT, space_after=18, line_spacing=1.15)


def s05_litsurvey(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Literature Survey", 5, PURPLE)

    add_text(slide, Inches(0.4), Inches(0.98), Inches(12.55), Inches(0.3),
             "Six topics shaped the design, and each lesson from past work led to a clear choice in AgeVision.",
             size=12.5, italic=True, color=MUTED)

    themes = [
        ("Guessing age with deep learning",
         "A 2024 NIST study found that single models stop improving at 4 to 7 years of error. "
         "So AgeVision uses two models together, namely MiVOLO v2 (about 3.65 yr) with an InsightFace backup.", PURPLE),
        ("Ageing a face with AI",
         "Past work shows keeping the person's identity is the hard part. AgeVision checks identity "
         "with FaceNet (score above 0.6) and offers three engines: SAM, Fast-AgingGAN and FADING.", ORANGE),
        ("Finding and recognising faces",
         "YOLOv8 finds faces first; RetinaFace is the backup. ArcFace codes help check identity, and "
         "OpenCV prepares both uploaded photos and live video.", GREEN),
        ("Reading emotions",
         "A ViT model from HuggingFace sorts each face into 7 emotions. It loads once and re-uses the "
         "faces already found, so it adds no extra work.", PURPLE),
        ("Clear results and two-way ageing",
         "Earlier studies asked for confidence scores and the ability to both age and de-age. AgeVision "
         "shows a confidence value and uses FADING for two-way ageing; visual heatmaps are planned next.", ORANGE),
        ("Web frameworks for AI apps",
         "Research recommends Django REST for AI apps. AgeVision pairs it with Angular 19 to keep the AI "
         "separate from the screens. Rival tools each do only part of the job.", GREEN),
    ]
    cw, ch, gx, gy = Inches(6.18), Inches(1.72), Inches(0.18), Inches(0.16)
    for i, (t, body, col) in enumerate(themes):
        col = (PURPLE, GREEN)[i % 2]
        x = Inches(0.4) + (i % 2) * (cw + gx)
        y = Inches(1.42) + (i // 2) * (ch + gy)
        card(slide, x, y, cw, ch, accent=col, border_pt=1.25)
        add_rect(slide, x, y, Inches(0.1), ch, col)
        add_text(slide, x + Inches(0.24), y + Inches(0.12), cw - Inches(0.42), Inches(0.32),
                 t, size=13, bold=True, color=col)
        add_text(slide, x + Inches(0.24), y + Inches(0.5), cw - Inches(0.42), ch - Inches(0.6),
                 body, size=11, color=TEXT, line_spacing=1.0)


def s06_proposed(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Proposed System", 6, GREEN)

    add_text(slide, Inches(0.4), Inches(0.98), Inches(12.55), Inches(0.32),
             "AgeVision brings everything into one free, self-hosted web app behind a single login.",
             size=12.5, italic=True, color=MUTED)

    section_label(slide, Inches(0.4), Inches(1.45), Inches(6.05), "What AgeVision offers", PURPLE)
    add_bullets(slide, Inches(0.4), Inches(2.0), Inches(6.05), Inches(4.5), [
        "High-accuracy age prediction from two models working together.",
        "Three engines to age or de-age a face, all keeping identity.",
        "Emotion reading on every face along with the age.",
        "Many photos at once (batch) and live camera input.",
        "Personal history, charts dashboard and settings.",
        "Admin panel with user management and live health checks.",
        "21 documented REST endpoints for outside apps.",
    ], size=12.5, color=TEXT, space_after=12, line_spacing=1.05)

    section_label(slide, Inches(6.75), Inches(1.45), Inches(6.2), "Main uses", GREEN)
    add_bullets(slide, Inches(6.75), Inches(2.0), Inches(6.2), Inches(4.5), [
        ("Forensics: ", "age progression to help find missing people."),
        ("Healthcare: ", "compare biological and chronological age."),
        ("Security: ", "age-gated access checks."),
        ("Research: ", "bundled MiVOLO v2 and SAM GAN training steps."),
        ("Bulk work: ", "process large case files through the batch module."),
    ], size=12.5, color=TEXT, space_after=15, line_spacing=1.06)


def s07_requirements(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "System Requirements", 7, GREEN)

    section_label(slide, Inches(0.4), Inches(1.0), Inches(6.05), "Hardware (Table 2.1)", PURPLE)
    hw = slide.shapes.add_table(6, 3, Inches(0.4), Inches(1.5), Inches(6.05), Inches(3.2)).table
    hw.columns[0].width = Inches(1.65); hw.columns[1].width = Inches(2.2); hw.columns[2].width = Inches(2.2)
    fill_table(hw, [
        ["Component", "Minimum", "Recommended"],
        ["Processor", "i5 8th Gen / Ryzen 5", "i7 / Ryzen 7 (8-core)"],
        ["RAM", "8 GB DDR4", "16 GB+ DDR4"],
        ["Storage", "50 GB HDD", "256 GB SSD"],
        ["GPU", "GTX 1060 (CUDA)", "RTX 3070 / Tesla T4"],
        ["Cloud GPU", "None", "Modal A10G / A100 (opt.)"],
    ], header_fill=PURPLE, hdr_pt=12, body_pt=11)

    section_label(slide, Inches(6.75), Inches(1.0), Inches(6.2), "Software (Table 2.2)", GREEN)
    sw = slide.shapes.add_table(9, 2, Inches(6.75), Inches(1.5), Inches(6.2), Inches(5.4)).table
    sw.columns[0].width = Inches(2.7); sw.columns[1].width = Inches(3.5)
    fill_table(sw, [
        ["Part", "Technology / Version"],
        ["Frontend (screens)", "Angular 19, Bootstrap 5.3, Chart.js 4, RxJS 7.8"],
        ["Backend (API)", "Python 3.11+, Django 5.2, DRF 3.16"],
        ["Main age model", "MiVOLO v2 (ViT + YOLOv8), ~3.65 MAE"],
        ["Backup age model", "InsightFace buffalo_l (RetinaFace + ArcFace)"],
        ["Ageing engines", "SAM GAN, Fast-AgingGAN, FADING diffusion"],
        ["Emotion / face finding", "ViT face-expression, YOLOv8 + RetinaFace"],
        ["Databases", "MongoDB 4.6 (pymongo) + SQLite (Django ORM)"],
        ["AI / security", "PyTorch 2.0+, diffusers 0.25+, JWT+Bcrypt+Fernet"],
    ], header_fill=GREEN, odd=RGBColor(0xED, 0xE8, 0xFC), even=RGBColor(0xF8, 0xF6, 0xFF),
       hdr_pt=12, body_pt=10.5)

    add_rect(slide, Inches(0.4), Inches(5.0), Inches(6.05), Inches(2.0), CARD2,
             border_color=BORDER, border_pt=1.0)
    add_text(slide, Inches(0.62), Inches(5.15), Inches(5.6), Inches(0.3),
             "Deployment notes", size=13, bold=True, color=ORANGE)
    add_bullets(slide, Inches(0.62), Inches(5.55), Inches(5.6), Inches(1.4), [
        "Runs on Windows 10 or Ubuntu 20.04+ (Ubuntu 22.04 for live use).",
        "Needs internet to download model files (HuggingFace, Modal).",
        "FADING ageing works best with a CUDA GPU or a Modal cloud GPU.",
    ], size=11.5, color=TEXT, space_after=8, line_spacing=1.0)


def s08_architecture(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Four-Tier System Design", 8, PURPLE)

    add_image_fit(slide, ss("fig_3_1_architecture.png"), Inches(0.4), Inches(1.1),
                  Inches(7.95), Inches(6.0), border=BORDER, valign="top")

    tiers = [
        ("Screens (Presentation)", GREEN,
         "What users see, made with Angular 19, Bootstrap 5, Chart.js and RxJS."),
        ("API (Business logic)", PURPLE,
         "Handles requests and security using Django 5.2 with DRF, JWT, CORS and 21 API calls."),
        ("AI / ML models", ORANGE,
         "Runs the AI models, including MiVOLO v2, InsightFace, SAM GAN, Fast-AgingGAN, FADING, ViT emotion and YOLOv8."),
        ("Storage (Data)", GREEN,
         "Keeps data in MongoDB 4.6 (6 sets), SQLite (logins) and an optional Modal cloud GPU."),
    ]
    rx, ry, rw, rh = Inches(8.55), Inches(1.1), Inches(4.4), Inches(1.42)
    for i, (t, col, body) in enumerate(tiers):
        col = (PURPLE, GREEN)[i % 2]
        card(slide, rx, ry, rw, rh, accent=col, border_pt=1.5)
        add_rect(slide, rx, ry, Inches(0.12), rh, col)
        add_text(slide, rx + Inches(0.28), ry + Inches(0.12), rw - Inches(0.45), Inches(0.34),
                 f"{i+1}.  {t}", size=13, bold=True, color=col)
        add_text(slide, rx + Inches(0.28), ry + Inches(0.54), rw - Inches(0.45), rh - Inches(0.62),
                 body, size=11, color=TEXT, line_spacing=1.0)
        ry += rh + Inches(0.1)


def s09_usecase(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Use Case Diagram", 9, PURPLE)

    add_image_fit(slide, ss("fig_2_1_use_case.png"), Inches(0.4), Inches(1.1),
                  Inches(7.6), Inches(5.95), border=BORDER, valign="top")

    card(slide, Inches(8.2), Inches(1.1), Inches(4.75), Inches(2.85), accent=PURPLE, border_pt=1.5)
    add_text(slide, Inches(8.42), Inches(1.24), Inches(4.3), Inches(0.3),
             "Four actors", size=14, bold=True, color=PURPLE)
    add_divider(slide, Inches(8.42), Inches(1.62), Inches(4.3), PURPLE, 1.25)
    add_bullets(slide, Inches(8.42), Inches(1.76), Inches(4.3), Inches(2.1), [
        ("User: ", "signs in and runs prediction, ageing, emotion and history."),
        ("Admin: ", "manages users and watches system health."),
        ("Platform: ", "the Angular client plus the Django backend."),
        ("External: ", "HuggingFace, YOLOv8 and Modal GPU."),
    ], size=11.5, color=TEXT, space_after=6, line_spacing=1.0)

    card(slide, Inches(8.2), Inches(4.1), Inches(4.75), Inches(2.9), accent=GREEN, border_pt=1.5)
    add_text(slide, Inches(8.42), Inches(4.24), Inches(4.3), Inches(0.3),
             "Main use cases", size=14, bold=True, color=GREEN)
    add_divider(slide, Inches(8.42), Inches(4.62), Inches(4.3), GREEN, 1.25)
    add_bullets(slide, Inches(8.42), Inches(4.76), Inches(4.3), Inches(2.1), [
        "Register, log in and reset password.",
        "Predict age for one, many or group photos.",
        "Age a face and read emotion.",
        "View history and charts.",
        "Admin: search, suspend and reinstate users.",
    ], size=11.5, color=TEXT, space_after=7, line_spacing=1.0)


def s10_dfd(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Data Flow Diagrams", 10, GREEN)

    add_text(slide, Inches(0.4), Inches(1.0), Inches(6.18), Inches(0.3),
             "Level 0: the whole system", size=13, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_image_fit(slide, ss("fig_2_2_dfd_l0.png"), Inches(0.4), Inches(1.4),
                  Inches(6.18), Inches(3.7), border=BORDER, valign="top")

    add_text(slide, Inches(6.75), Inches(1.0), Inches(6.2), Inches(0.3),
             "Level 1: eight sub-processes", size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_image_fit(slide, ss("fig_2_3_dfd_l1.png"), Inches(6.75), Inches(1.4),
                  Inches(6.2), Inches(3.7), border=BORDER, valign="top")

    add_rect(slide, Inches(0.4), Inches(5.45), Inches(12.55), Inches(1.55), CARD,
             border_color=PURPLE, border_pt=1.0)
    add_text(slide, Inches(0.62), Inches(5.58), Inches(12.1), Inches(1.3),
             "Level 0 shows AgeVision as one process that takes photo uploads, batches, live camera frames "
             "and logins from the User and Admin, and returns ages, aged images and admin reports. "
             "Level 1 breaks it into eight steps: login, media upload, face finding, age and emotion, "
             "multi-engine ageing, history, analytics and admin. Each step owns its own MongoDB collection "
             "or SQLite table.",
             size=12, color=TEXT, line_spacing=1.12, anchor=MSO_ANCHOR.MIDDLE)


def s11_database(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Database Design", 11, PURPLE)

    add_image_fit(slide, ss("fig_3_3_er_diagram.png"), Inches(0.4), Inches(1.1),
                  Inches(7.0), Inches(5.95), border=BORDER, valign="top")

    card(slide, Inches(7.6), Inches(1.1), Inches(5.35), Inches(5.9), accent=PURPLE, border_pt=1.5)
    add_text(slide, Inches(7.82), Inches(1.26), Inches(4.9), Inches(0.3),
             "Hybrid storage", size=14.5, bold=True, color=PURPLE)
    add_divider(slide, Inches(7.82), Inches(1.66), Inches(4.9), PURPLE, 1.25)
    add_bullets(slide, Inches(7.82), Inches(1.82), Inches(4.9), Inches(2.3), [
        ("SQLite (Django ORM): ", "users, sessions and tokens."),
        ("MongoDB 4.6 (PyMongo): ", "all AI results in six collections."),
        ("One User has ", "many Predictions and Progressions, exactly one Settings, and many Batch Jobs."),
        ("Each Batch Job ", "groups many Predictions; each Prediction can hold many faces."),
    ], size=12, color=TEXT, space_after=10, line_spacing=1.05)

    add_text(slide, Inches(7.82), Inches(5.0), Inches(4.9), Inches(0.3),
             "Six MongoDB collections", size=12.5, bold=True, color=GREEN)
    add_bullets(slide, Inches(7.82), Inches(5.36), Inches(4.9), Inches(1.5), [
        "predictions, progressions, user_settings,",
        "batch_predictions, password_resets, history.",
    ], size=12, color=TEXT, space_after=4, line_spacing=1.05)


def s12_modules(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "List of Modules", 12, GREEN)

    add_text(slide, Inches(0.4), Inches(0.96), Inches(12.55), Inches(0.32),
             "Thirteen features, all built on Django and shown through the Angular 19 app.",
             size=12.5, italic=True, color=MUTED)

    modules = [
        ("Dashboard", "Totals, recent results and quick actions."),
        ("Login & Sign-up", "Register, log in and reset password (JWT)."),
        ("Single-Face Age", "Find one face and guess its age."),
        ("Group Photo", "Age every face in a group photo."),
        ("Many Photos (Batch)", "Upload and process many images together."),
        ("Emotion Reading", "Label each face with 1 of 7 emotions."),
        ("Face Ageing", "Age or de-age a face with three engines."),
        ("Live Progress", "Show a live progress bar for slow jobs (SSE)."),
        ("History", "Saved predictions, progressions and batches."),
        ("Charts Dashboard", "Counts, age spread, gender and emotion charts."),
        ("Settings", "Theme, preferred model and engine, alerts."),
        ("Live Camera", "Read webcam frames and overlay age live."),
        ("Admin Panel", "User management and live system health."),
    ]
    col_w, gx = Inches(6.18), Inches(0.16)
    x_cols = [Inches(0.4), Inches(0.4) + col_w + gx]
    y0, pitch, chh = Inches(1.42), Inches(0.66), Inches(0.58)
    badge_off = Inches(0.14)

    def mod_card(x, y, w, idx, name, desc):
        col = (PURPLE, GREEN)[idx % 2]
        card(slide, x, y, w, chh, accent=col, border_pt=1.0)
        add_rect(slide, x + Inches(0.12), y + badge_off, Inches(0.3), Inches(0.3), col, rounded=True)
        add_text(slide, x + Inches(0.12), y + badge_off, Inches(0.3), Inches(0.3),
                 str(idx + 1), size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.55), y + Inches(0.05), w - Inches(0.66), Inches(0.26),
                 name, size=11.5, bold=True, color=TEXT)
        add_text(slide, x + Inches(0.55), y + Inches(0.3), w - Inches(0.66), Inches(0.24),
                 desc, size=9.5, color=MUTED)

    for i in range(12):
        mod_card(x_cols[i % 2], y0 + (i // 2) * pitch, col_w, i, *modules[i])
    mod_card(x_cols[0], y0 + 6 * pitch, Inches(12.52), 12, *modules[12])


def s13_prediction(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Age Prediction Using Two Models", 13, PURPLE)

    section_label(slide, Inches(0.4), Inches(1.0), Inches(5.85), "MiVOLO v2 + InsightFace", GREEN)
    add_bullets(slide, Inches(0.4), Inches(1.55), Inches(5.85), Inches(3.4), [
        ("Main model: MiVOLO v2", " looks at both the face and the body."),
        ("Backup model: InsightFace", " is used only if needed (about 8.5 yr error)."),
        ("Combined result", " is 70% MiVOLO plus 30% InsightFace."),
        ("Finding faces", " uses YOLOv8 first, then RetinaFace as backup."),
        ("Group photos", " give the age of every face, with gender and emotion."),
        ("Confidence", " shows how sure the result is, per face."),
        ("Live camera", " reads a new frame every 1.5 seconds."),
    ], size=12.5, color=TEXT, space_after=7, line_spacing=1.0)

    add_rect(slide, Inches(0.4), Inches(5.0), Inches(5.85), Inches(1.0), CHIP,
             border_color=PURPLE, border_pt=1.25, rounded=True)
    add_text(slide, Inches(0.4), Inches(5.12), Inches(5.85), Inches(0.4),
             "Average error  about  3.65 years", size=20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.4), Inches(5.56), Inches(5.85), Inches(0.34),
             "About 48 to 54% fewer errors than How-Old.net or Face++ (6 to 8 yr)", size=11.5,
             color=MUTED, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.4), Inches(6.18), Inches(5.85), Inches(0.9),
             "Emotion: a ViT model sorts each face into 7 feelings (angry, disgust, fear, happy, "
             "neutral, sad, surprise), re-using the faces already found.",
             size=11, italic=True, color=MUTED, line_spacing=1.0)

    add_image_fit(slide, ss("fig_4_3_mivolo_pipeline.png"), Inches(6.55), Inches(1.05),
                  Inches(6.4), Inches(2.85), border=BORDER)
    add_image_fit(slide, ss("fig_4_2_age_prediction.png"), Inches(6.55), Inches(4.05),
                  Inches(6.4), Inches(2.95), border=PURPLE)


def s14_progression(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Age Progression Using Three Engines", 14, ORANGE)

    add_text(slide, Inches(0.4), Inches(0.98), Inches(12.55), Inches(0.32),
             "Pick one of three engines. Every result must pass a FaceNet identity check (score above 0.6).",
             size=12.5, italic=True, color=MUTED)

    engines = [
        ("SAM GAN", PURPLE, "MAIN",
         "Turns the photo into StyleGAN2 code, then shifts it toward the target age. "
         "Tuned for Indian faces. Takes 3 to 8 seconds."),
        ("Fast-AgingGAN", GREEN, "RUNS ON CPU",
         "A small CycleGAN (about 11 MB). Works on a normal computer with no GPU, which is good for low-power setups."),
        ("FADING Diffusion", ORANGE, "BEST QUALITY",
         "Uses Stable Diffusion to age or de-age a face. Highest quality; can use a Modal cloud GPU."),
    ]
    cw, gap = Inches(4.05), Inches(0.2)
    for i, (name, col, badge, body) in enumerate(engines):
        x = Inches(0.4) + i * (cw + gap)
        card(slide, x, Inches(1.42), cw, Inches(2.5), accent=col, border_pt=1.5)
        add_text(slide, x + Inches(0.2), Inches(1.56), cw - Inches(0.4), Inches(0.4),
                 name, size=16, bold=True, color=col)
        add_rect(slide, x + Inches(0.2), Inches(2.0), Inches(1.85), Inches(0.34), CHIP,
                 border_color=col, border_pt=1.0, rounded=True)
        add_text(slide, x + Inches(0.2), Inches(2.0), Inches(1.85), Inches(0.34), badge,
                 size=10.5, bold=True, color=col, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.2), Inches(2.46), cw - Inches(0.4), Inches(1.4),
                 body, size=11.5, color=TEXT, line_spacing=1.02)

    add_image_fit(slide, ss("fig_4_4_sam_gan_flow.png"), Inches(0.4), Inches(4.1),
                  Inches(7.7), Inches(2.9), border=BORDER)

    mx, mw, my = Inches(8.4), Inches(4.55), Inches(4.1)
    for big, small in [
        ("> 95%", "of faces keep the same identity"),
        ("3 to 8 s", "to age a face with SAM GAN"),
        ("Two-way", "FADING can age and de-age a face"),
    ]:
        add_rect(slide, mx, my, mw, Inches(0.84), CHIP, border_color=PURPLE,
                 border_pt=1.0, rounded=True)
        add_text(slide, mx + Inches(0.2), my, Inches(1.85), Inches(0.84), big,
                 size=18, bold=True, color=PURPLE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, mx + Inches(2.1), my, mw - Inches(2.3), Inches(0.84), small,
                 size=11.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        my += Inches(0.99)


def s15_detection(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Detection & Input Modules", 15, PURPLE)
    demo_grid(slide, [
        (ss("fig_4_1_auth_flow.png"), "Secure login, sign-up and password reset", PURPLE),
        (ss("fig_4_11_age_group_emotion_prediction.png"), "Group photo with age, gender and emotion", GREEN),
        (ss("fig_4_7_batch_predict.png"), "Batch: many photos processed together", ORANGE),
        (ss("fig_4_9_camera.png"), "Live camera with age overlay", GREEN),
    ])


def s16_analytics(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "History, Analytics & Admin", 16, GREEN)
    demo_grid(slide, [
        (ss("fig_4_8_dashboard.png"), "Home dashboard with totals and recent results", PURPLE),
        (ss("fig_4_12_analysis_dashboard_1.png"), "Charts for age, gender and emotion", GREEN),
        (ss("fig_4_15_analysis_history.png"), "History with thumbnails, download and delete", PURPLE),
        (ss("fig_4_14_admin_panel.png"), "Admin: user management and system health", GREEN),
    ])


def s17_datasets(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Datasets & Algorithms", 17, GREEN)

    section_label(slide, Inches(0.4), Inches(1.0), Inches(6.05), "Datasets (Table 4.1)", PURPLE)
    dt = slide.shapes.add_table(5, 3, Inches(0.4), Inches(1.5), Inches(6.05), Inches(2.35)).table
    dt.columns[0].width = Inches(1.95); dt.columns[1].width = Inches(1.5); dt.columns[2].width = Inches(2.6)
    fill_table(dt, [
        ["Dataset", "Size", "Used For"],
        ["UTKFace", "23,000+", "Train + test MiVOLO v2"],
        ["IMDB-WIKI", "500,000+", "Extra training for MiVOLO v2"],
        ["MORPH", "55,000", "Testing only"],
        ["FFHQ", "70,000", "Train SAM GAN + FADING"],
    ], header_fill=PURPLE, hdr_pt=12, body_pt=11)

    add_text(slide, Inches(0.4), Inches(4.1), Inches(6.05), Inches(2.9),
             "An Indian-face subset (filtered with YOLOv8 confidence of at least 0.8) was added to FFHQ "
             "to improve results on South-Asian faces. MORPH is kept aside only for testing, never for training.",
             size=12, color=TEXT, line_spacing=1.15)

    card(slide, Inches(6.75), Inches(1.0), Inches(6.2), Inches(6.0), accent=GREEN, border_pt=1.5)
    add_text(slide, Inches(6.97), Inches(1.16), Inches(5.7), Inches(0.3),
             "How the algorithms work", size=14.5, bold=True, color=GREEN)
    add_divider(slide, Inches(6.97), Inches(1.56), Inches(5.7), GREEN, 1.25)
    add_bullets(slide, Inches(6.97), Inches(1.72), Inches(5.7), Inches(5.1), [
        ("MiVOLO v2: ", "a two-head Vision Transformer that also reads body cues from YOLOv8. The final age is 70% MiVOLO plus 30% InsightFace."),
        ("SAM GAN: ", "a pSp encoder turns the photo into StyleGAN2 W+ code, then moves it toward the target-age direction."),
        ("FADING: ", "null-text inversion rebuilds the photo's diffusion path, then classifier-free guidance steers it to the target age, so it can age and de-age."),
        ("Identity check: ", "FaceNet cosine similarity between the old and new face must stay above 0.6."),
    ], size=12, color=TEXT, space_after=12, line_spacing=1.06)


def s18_testing(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Testing & Validation", 18, PURPLE)

    section_label(slide, Inches(0.4), Inches(1.0), Inches(6.2), "What we checked in testing", GREEN)
    tt = slide.shapes.add_table(9, 3, Inches(0.4), Inches(1.5), Inches(6.2), Inches(4.55)).table
    tt.columns[0].width = Inches(2.9); tt.columns[1].width = Inches(2.05); tt.columns[2].width = Inches(1.25)
    fill_table(tt, [
        ["Area", "What it checks", "Tests"],
        ["Login (TC-A)", "Sign-up, login, JWT, reset", "10"],
        ["Prediction (TC-P)", "Single, group, camera", "9"],
        ["Batch (TC-P)", "Many photos and saving", "2"],
        ["Ageing (TC-G)", "SAM / Fast / FADING + progress", "11"],
        ["History (TC-H)", "List, delete, user privacy", "4"],
        ["Charts (TC-N)", "Totals and chart data", "3"],
        ["Settings (TC-S)", "Theme, model, checks", "4"],
        ["Admin (TC-AD)", "Stats, suspend, health", "7"],
    ], header_fill=GREEN, odd=RGBColor(0xED, 0xE8, 0xFC), even=RGBColor(0xF8, 0xF6, 0xFF),
       hdr_pt=12, body_pt=11)

    section_label(slide, Inches(6.75), Inches(1.0), Inches(6.2), "Speed (Table 4.5)", ORANGE)
    tc = slide.shapes.add_table(5, 2, Inches(6.75), Inches(1.5), Inches(6.2), Inches(2.4)).table
    tc.columns[0].width = Inches(3.9); tc.columns[1].width = Inches(2.3)
    fill_table(tc, [
        ["Step", "Time grows with"],
        ["Finding faces (YOLOv8)", "O(n) by pixels"],
        ["Guessing age (2 models)", "O(f) by faces"],
        ["SAM GAN ageing", "O(1) per image"],
        ["FADING ageing", "O(T) by steps"],
    ], header_fill=PURPLE, odd=RGBColor(0xED, 0xE8, 0xFC), even=RGBColor(0xF8, 0xF6, 0xFF),
       hdr_pt=12, body_pt=11)

    stat_chips(slide, Inches(4.35), [
        ("50", "tests, all passed", PURPLE),
        ("13", "features built", GREEN),
        ("21", "REST endpoints", PURPLE),
    ], x0=Inches(6.75), total_w=Inches(6.2))

    add_text(slide, Inches(6.75), Inches(5.5), Inches(6.2), Inches(1.4),
             "Every module was checked against a written test catalogue of 50 cases covering normal use, "
             "wrong input and edge cases. All 50 passed.",
             size=12, color=TEXT, line_spacing=1.15)


def s19_results(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Results: How AgeVision Compares", 19, GREEN)

    tbl = slide.shapes.add_table(8, 6, Inches(0.4), Inches(1.05), Inches(8.0), Inches(5.6)).table
    widths = [2.0, 1.35, 1.2, 1.0, 1.25, 1.2]
    for i, wd in enumerate(widths):
        tbl.columns[i].width = Inches(wd)
    fill_table(tbl, [
        ["Feature", "AgeVision", "How-Old", "Face++", "AgingBooth", "DeepFace"],
        ["Age error (MAE)", "~3.65 yr", "6 to 8 yr", "5 to 7 yr", "N/A", "~5.4 yr"],
        ["Ageing engines", "3 (GAN+Diff)", "None", "None", "1 (GAN)", "None"],
        ["Many photos at once", "Yes", "No", "Paid", "No", "No"],
        ["Emotion reading", "7-class ViT", "No", "Yes", "No", "Yes"],
        ["Keeps identity", "FaceNet ≥0.6", "No", "No", "No", "No"],
        ["Free / open-source", "Yes", "Yes", "Paid", "Paid", "Yes"],
        ["Admin + self-hosted", "Yes", "No", "Paid", "No", "Partly"],
    ], header_fill=GREEN, odd=RGBColor(0xED, 0xE8, 0xFC), even=RGBColor(0xF8, 0xF6, 0xFF),
       hdr_pt=11.5, body_pt=10.5, last_fill=None, align_first_left=True)
    add_text(slide, Inches(0.4), Inches(6.72), Inches(8.0), Inches(0.24),
             "Table 4.6 shows that AgeVision is the only free, self-hosted tool that does all of this in one place.",
             size=10, italic=True, color=MUTED)

    add_image_fit(slide, ss("Figure_4_6_Comparison.png"), Inches(8.55), Inches(1.05),
                  Inches(4.4), Inches(4.3), border=BORDER, valign="top")

    add_rect(slide, Inches(8.55), Inches(5.5), Inches(4.4), Inches(1.5), CHIP,
             border_color=PURPLE, border_pt=1.25, rounded=True)
    add_text(slide, Inches(8.75), Inches(5.62), Inches(4.0), Inches(0.32),
             "Headline result", size=13, bold=True, color=PURPLE)
    add_text(slide, Inches(8.75), Inches(5.98), Inches(4.0), Inches(0.95),
             "About 3.65 years of error, which is roughly 48 to 54% fewer mistakes than paid tools, "
             "and over 95% of aged faces keep the same identity.",
             size=11.5, color=TEXT, line_spacing=1.0)


def s20_conclusion(prs, layout):
    slide = prs.slides.add_slide(layout); set_bg(slide)
    header(slide, "Conclusion & Future Work", 20, ORANGE)

    card(slide, Inches(0.4), Inches(1.05), Inches(6.15), Inches(4.5), accent=GREEN, border_pt=1.5)
    add_text(slide, Inches(0.62), Inches(1.2), Inches(5.7), Inches(0.34),
             "What we achieved", size=14.5, bold=True, color=GREEN)
    add_divider(slide, Inches(0.62), Inches(1.58), Inches(5.7), GREEN, 1.5)
    add_bullets(slide, Inches(0.62), Inches(1.74), Inches(5.7), Inches(3.7), [
        "Two-model age guessing gives about 3.65 years of error (48 to 54% better than paid tools).",
        "Three ageing engines, with over 95% keeping the same identity.",
        "FADING can both age and de-age a face realistically.",
        "13 features, 21 API calls and 6 MongoDB sets, all login-protected.",
        "Charts dashboard plus an admin-only control page.",
        "Re-usable training steps for MiVOLO v2 and SAM GAN included.",
    ], size=11.5, color=TEXT, space_after=11, line_spacing=1.04)

    card(slide, Inches(6.75), Inches(1.05), Inches(6.2), Inches(4.5), accent=ORANGE, border_pt=1.5)
    add_text(slide, Inches(6.97), Inches(1.2), Inches(5.7), Inches(0.34),
             "What's next", size=14.5, bold=True, color=ORANGE)
    add_divider(slide, Inches(6.97), Inches(1.58), Inches(5.7), ORANGE, 1.5)
    add_bullets(slide, Inches(6.97), Inches(1.74), Inches(5.7), Inches(3.7), [
        "Host fully in the cloud (AWS / GCP GPU) with a job queue for slow tasks.",
        "Add a PDF download for each result and a mobile app.",
        "Add visual heatmaps and Tamil / Hindi screens.",
        "Connect to police case systems and speed up admin charts.",
        "Auto-retrain the models as new photos arrive.",
        "Bundle model files for offline (no-internet) use.",
    ], size=11.5, color=TEXT, space_after=13, line_spacing=1.04)

    add_rect(slide, Inches(0.4), Inches(5.7), Inches(12.55), Inches(0.66), CARD2,
             border_color=BORDER, border_pt=1.0)
    add_text(slide, Inches(0.6), Inches(5.7), Inches(12.15), Inches(0.66),
             "Key references:  MiVOLO (2025)  •  NIST FATE 8525 (2024)  •  Null-text Inversion, CVPR 2023  "
             "•  SAM: Only a Matter of Style  •  Django REST Framework",
             size=10.5, italic=True, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(slide, Inches(0.4), Inches(6.46), Inches(12.55), Inches(0.6), CHIP,
             border_color=PURPLE, border_pt=1.0, rounded=True)
    add_text(slide, Inches(0.4), Inches(6.46), Inches(12.55), Inches(0.6),
             "Thank You and Open to Questions", size=18, bold=True, color=PURPLE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ═════════════════════════════════════════════════════════════════════════ #
#                                   MAIN                                     #
# ═════════════════════════════════════════════════════════════════════════ #

def ensure_app_logo():
    if os.path.exists(APP_LOGO):
        return
    svg = """<svg viewBox="0 0 120 80" xmlns="http://www.w3.org/2000/svg">
  <path d="M60 10 C30 10, 5 40, 5 40 C5 40, 30 70, 60 70 C90 70, 115 40, 115 40 C115 40, 90 10, 60 10Z" fill="none" stroke="#1a3a5c" stroke-width="3.5"/>
  <line x1="5" y1="40" x2="0" y2="40" stroke="#c8922a" stroke-width="2.5"/>
  <line x1="10" y1="28" x2="2" y2="22" stroke="#c8922a" stroke-width="2"/>
  <line x1="10" y1="52" x2="2" y2="58" stroke="#c8922a" stroke-width="2"/>
  <line x1="115" y1="40" x2="120" y2="40" stroke="#c8922a" stroke-width="2.5"/>
  <line x1="110" y1="28" x2="118" y2="22" stroke="#c8922a" stroke-width="2"/>
  <line x1="110" y1="52" x2="118" y2="58" stroke="#c8922a" stroke-width="2"/>
  <circle cx="60" cy="40" r="22" fill="none" stroke="#c8922a" stroke-width="3" stroke-dasharray="18 6"/>
  <circle cx="60" cy="40" r="15" fill="none" stroke="#1a3a5c" stroke-width="2.5"/>
  <circle cx="60" cy="40" r="9" fill="none" stroke="#1a3a5c" stroke-width="2"/>
  <circle cx="60" cy="40" r="4" fill="#1a3a5c"/>
  <circle cx="57" cy="37" r="1.5" fill="#fff" opacity="0.8"/>
</svg>"""
    try:
        import fitz
        tmp = os.path.join(ROOT, "_logo_tmp.svg")
        with open(tmp, "w") as f:
            f.write(svg)
        doc = fitz.open(tmp)
        doc[0].get_pixmap(matrix=fitz.Matrix(9, 9), alpha=True).save(APP_LOGO)
        doc.close()
        os.remove(tmp)
    except Exception as e:
        print(f"  (could not render logo: {e})")


def build():
    ensure_app_logo()
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    builders = [
        ("Title", s01_title),
        ("Abstract", s02_abstract),
        ("Introduction & Problem", s03_intro),
        ("Objectives & Scope", s04_objectives),
        ("Literature Survey", s05_litsurvey),
        ("Proposed System", s06_proposed),
        ("System Requirements", s07_requirements),
        ("System Architecture", s08_architecture),
        ("Use Case Diagram", s09_usecase),
        ("Data Flow Diagrams", s10_dfd),
        ("Database Design", s11_database),
        ("List of Modules", s12_modules),
        ("Age Prediction", s13_prediction),
        ("Age Progression", s14_progression),
        ("Detection & Input", s15_detection),
        ("History, Analytics & Admin", s16_analytics),
        ("Datasets & Algorithms", s17_datasets),
        ("Testing & Validation", s18_testing),
        ("Results & Comparison", s19_results),
        ("Conclusion & Future Work", s20_conclusion),
    ]
    print("Building slides...")
    for i, (name, fn) in enumerate(builders, 1):
        fn(prs, blank)
        print(f"  [{i:2d}/{TOTAL}] {name}")

    out = OUT
    try:
        prs.save(out)
    except PermissionError:
        out = OUT.replace(".pptx", "_new.pptx")
        prs.save(out)
        print(f"WARNING: original was open. Saved to {out}")
    print(f"\nSaved PPTX: {out}  ({len(prs.slides)} slides)")
    return out


def to_pdf(pptx_path):
    try:
        import win32com.client
    except ImportError:
        print("win32com not available, skipping PDF export.")
        return None
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    deck = None
    try:
        deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
        deck.SaveAs(os.path.abspath(PDF), 32)  # 32 = ppSaveAsPDF
        print(f"Saved PDF:  {PDF}")
    finally:
        if deck is not None:
            deck.Close()
        powerpoint.Quit()
    return PDF


if __name__ == "__main__":
    pptx = build()
    to_pdf(pptx)
