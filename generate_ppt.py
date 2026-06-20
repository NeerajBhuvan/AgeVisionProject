"""Generate AgeVision_VivaPPT.pptx — 10-slide viva-voce presentation.

White background, purple/green/orange accents, clean card layouts.

Run: python generate_ppt.py
"""
from __future__ import annotations

import os
import subprocess
import sys

# ── auto-install dependency ──────────────────────────────────────────────────
try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

# ── paths ────────────────────────────────────────────────────────────────────
ROOT   = os.path.dirname(os.path.abspath(__file__))
SS_DIR = os.path.join(ROOT, "screenshots")
LOGO   = os.path.join(ROOT, "age_vision_logo.png")
OUT    = os.path.join(ROOT, "AgeVision_VivaPPT.pptx")

def ss(name: str) -> str:
    return os.path.join(SS_DIR, name)

# ── slide dimensions (16:9 widescreen) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── colour palette — white/light theme ───────────────────────────────────────
BG     = RGBColor(0xFF, 0xFF, 0xFF)   # white background
CARD   = RGBColor(0xF4, 0xF1, 0xFC)   # very light purple-tinted card
BORDER = RGBColor(0xD4, 0xC8, 0xF5)   # soft purple border
PURPLE = RGBColor(0x7F, 0x5A, 0xF0)   # primary accent
GREEN  = RGBColor(0x2C, 0xB6, 0x7D)   # secondary accent
ORANGE = RGBColor(0xFF, 0x89, 0x06)   # tertiary accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFE)   # text on dark/coloured shapes
MUTED  = RGBColor(0x5C, 0x58, 0x80)   # captions/secondary text on white
TEXT   = RGBColor(0x1A, 0x16, 0x25)   # body text on white background
HDR_END = RGBColor(0x5A, 0x38, 0xC8)  # header gradient end (deep purple)

FONT = "Sora"


# ═══════════════════════════════════════════════════════════════════════════ #
#                              HELPER FUNCTIONS                               #
# ═══════════════════════════════════════════════════════════════════════════ #

def _rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color: RGBColor, *,
             border_color: RGBColor | None = None, border_pt: float = 0.75):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1 = RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def add_gradient_rect(slide, l, t, w, h,
                      c_start: RGBColor = PURPLE,
                      c_end: RGBColor = HDR_END,
                      angle: int = 0):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:solidFill")):
        spPr.remove(old)
    grad_xml = (
        f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{_rgb_hex(c_start)}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{_rgb_hex(c_end)}"/></a:gs>'
        f'</a:gsLst>'
        f'<a:lin ang="{angle}" scaled="0"/>'
        f'</a:gradFill>'
    )
    spPr.insert(0, etree.fromstring(grad_xml))
    return shape


def add_textbox(slide, l, t, w, h, text: str, *,
                font_size: float = 14,
                bold: bool = False,
                italic: bool = False,
                color: RGBColor = TEXT,
                align: PP_ALIGN = PP_ALIGN.LEFT,
                font_name: str = FONT,
                word_wrap: bool = True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def add_divider(slide, l, t, w, color: RGBColor = PURPLE, thickness_pt: float = 2):
    return add_rect(slide, l, t, w, Pt(thickness_pt), color)


def add_image(slide, path: str, l, t, w, h=None):
    if not os.path.exists(path):
        shape = add_rect(slide, l, t, w, h or Inches(2), CARD, border_color=BORDER)
        add_textbox(slide, l, t, w, h or Inches(2), os.path.basename(path),
                    font_size=9, color=MUTED, align=PP_ALIGN.CENTER)
        return shape
    kwargs = {"width": w}
    if h:
        kwargs["height"] = h
    return slide.shapes.add_picture(path, l, t, **kwargs)


def add_slide_header(slide, title: str, slide_num: int,
                     accent: RGBColor = PURPLE) -> None:
    # header bar: accent → deep purple gradient
    add_gradient_rect(slide, 0, 0, W, Inches(0.88), accent, HDR_END, angle=0)
    add_textbox(slide, Inches(0.45), Inches(0.08), Inches(11.6), Inches(0.72),
                title, font_size=24, bold=True, color=WHITE)
    add_rect(slide, Inches(12.55), Inches(0.18), Inches(0.55), Inches(0.55), accent)
    add_textbox(slide, Inches(12.55), Inches(0.18), Inches(0.55), Inches(0.55),
                f"{slide_num:02d}", font_size=13, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(0.4), Inches(0.90), Inches(12.53), accent, 2)
    # footer bar
    add_rect(slide, 0, Inches(7.38), W, Inches(0.12), accent)
    add_textbox(slide, Inches(0.4), Inches(7.36), Inches(10), Inches(0.14),
                "AgeVision  •  MCA IV Project  •  Anna University CDE  •  April 2026",
                font_size=8, color=MUTED)


def add_card(slide, l, t, w, h, *,
             title: str = "",
             title_color: RGBColor = PURPLE,
             items: list[str] | None = None,
             body_size: float = 13,
             padding: float = 0.18):
    add_rect(slide, l, t, w, h, CARD, border_color=BORDER, border_pt=0.75)
    y = t + Inches(padding)
    if title:
        add_textbox(slide, l + Inches(padding), y, w - Inches(2 * padding), Inches(0.35),
                    title, font_size=13, bold=True, color=title_color)
        add_divider(slide, l + Inches(padding), y + Inches(0.36),
                    w - Inches(2 * padding), PURPLE, 1.5)
        y += Inches(0.52)
    if items:
        for item in items:
            add_textbox(slide, l + Inches(padding), y,
                        w - Inches(2 * padding), Inches(0.38),
                        f"• {item}", font_size=body_size, color=TEXT)
            y += Inches(0.38)


def _set_cell_border(cell, color: RGBColor = BORDER):
    hex_color = _rgb_hex(color)
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for side in ("lnL", "lnR", "lnT", "lnB"):
        NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ln = etree.SubElement(tcPr, f"{{{NS}}}{side}")
        ln.set("w", "12700")
        solidFill = etree.SubElement(ln, f"{{{NS}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{NS}}}srgbClr")
        srgbClr.set("val", hex_color)


def style_table(tbl, header_fill: RGBColor, data_odd: RGBColor, data_even: RGBColor,
                hdr_text: RGBColor = WHITE, body_text: RGBColor = TEXT,
                hdr_pt: float = 13, body_pt: float = 12,
                last_row_fill: RGBColor | None = None,
                last_row_text: RGBColor = WHITE):
    nrows = len(tbl.rows)
    for ri, row in enumerate(tbl.rows):
        is_last = (ri == nrows - 1) and last_row_fill is not None
        for ci, cell in enumerate(row.cells):
            fill = cell.fill
            fill.solid()
            if ri == 0:
                fill.fore_color.rgb = header_fill
            elif is_last:
                fill.fore_color.rgb = last_row_fill
            elif ri % 2 == 1:
                fill.fore_color.rgb = data_odd
            else:
                fill.fore_color.rgb = data_even
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.04)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.name = FONT
                    run.font.size = Pt(hdr_pt if ri == 0 else body_pt)
                    run.font.bold = (ri == 0 or is_last)
                    if ri == 0:
                        run.font.color.rgb = hdr_text
                    elif is_last:
                        run.font.color.rgb = last_row_text
                    else:
                        run.font.color.rgb = body_text
            _set_cell_border(cell)


def _add_table_cell(tbl, row: int, col: int, text: str) -> None:
    cell = tbl.cell(row, col)
    cell.text_frame.text = text
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.name = FONT


# ═══════════════════════════════════════════════════════════════════════════ #
#                            SLIDE BUILDERS                                   #
# ═══════════════════════════════════════════════════════════════════════════ #

def build_slide_01_title(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)

    # top & bottom accent bars
    add_gradient_rect(slide, 0, 0, W, Inches(0.12), PURPLE, HDR_END, angle=0)
    add_gradient_rect(slide, 0, Inches(7.38), W, Inches(0.12), HDR_END, PURPLE, angle=0)

    # light purple tint behind content area
    add_rect(slide, 0, Inches(1.65), W, Inches(4.2), RGBColor(0xF9, 0xF7, 0xFF))

    # logo
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(5.4), Inches(0.22), width=Inches(2.5))

    # title
    add_textbox(slide, Inches(0.5), Inches(1.75), Inches(12.33), Inches(1.2),
                "AgeVision", font_size=56, bold=True, color=PURPLE,
                align=PP_ALIGN.CENTER)

    # subtitle
    add_textbox(slide, Inches(0.5), Inches(2.9), Inches(12.33), Inches(0.7),
                "AI Based Age Prediction and Age Progression System Using Facial Images",
                font_size=18, color=TEXT, align=PP_ALIGN.CENTER)

    # accent divider
    add_divider(slide, Inches(4.2), Inches(3.7), Inches(4.93), PURPLE, 2.5)

    # author block
    for y, text, col, sz in [
        (Inches(3.9),  "Neeraj Bhuvan M  |  Roll No. 2435MCA0011  |  Reg. No. 67224100038",
         GREEN, 15),
        (Inches(4.3),  "MCA IV  •  Anna University, Centre for Distance Education",
         MUTED, 14),
        (Inches(4.65), "Guide: Dr. P. Geetha, Associate Professor, IST Dept., CEG",
         MUTED, 13),
        (Inches(5.0),  "April 2026",
         ORANGE, 14),
    ]:
        add_textbox(slide, Inches(0.5), y, Inches(12.33), Inches(0.42),
                    text, font_size=sz, color=col, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.3),
                "Viva-Voce Presentation", font_size=10, color=MUTED,
                align=PP_ALIGN.CENTER)


def build_slide_02_problem(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Problem Statement & Motivation", 2, PURPLE)

    add_card(slide, Inches(0.35), Inches(1.05), Inches(6.0), Inches(6.0),
             title="The Gap in Existing Solutions", title_color=PURPLE,
             items=[
                 "No free, unified platform for age prediction + progression",
                 "Commercial APIs (Face++, AWS Rekognition) are paid & privacy-invasive",
                 "Open-source tools need developer expertise to deploy",
                 "Age progression lacks identity-preserving realism",
                 "No system supports batch forensic-scale processing",
                 "Missing emotion + demographic metadata alongside age",
             ], body_size=13)

    add_card(slide, Inches(6.6), Inches(1.05), Inches(6.4), Inches(6.0),
             title="Why It Matters", title_color=GREEN,
             items=[
                 "Forensics: missing-persons & cold-case age simulation",
                 "Healthcare: dermatological & geriatric age assessment",
                 "Entertainment: VFX de-aging / aging without specialist tools",
                 "Security: biometric age-gate & facial liveness verification",
                 "Academic: benchmark platform for generative AI research",
                 "Insurance / HR: age-group analytics at scale",
             ], body_size=13)

    # bottom highlight bar — light orange tint
    add_rect(slide, Inches(0.35), Inches(7.08), Inches(12.63), Inches(0.22),
             RGBColor(0xFF, 0xF0, 0xD0))
    add_textbox(slide, Inches(0.35), Inches(7.09), Inches(12.63), Inches(0.2),
                "AgeVision bridges the gap: one platform, multiple engines, no cost, no data leaves your machine.",
                font_size=10, color=ORANGE, italic=True, align=PP_ALIGN.CENTER)


def build_slide_03_objectives(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Core Objectives", 3, GREEN)

    objectives = [
        ("01", "Develop Angular 19 + Django 5.2 web platform with single, batch & real-time camera modes"),
        ("02", "Implement ensemble age prediction: MiVOLO v2 primary (MAE ≤ 3.65 yr) + InsightFace fallback"),
        ("03", "Build 4-engine age progression cascade: SAM Indian → SAM FFHQ → Fast-AgingGAN → FADING"),
        ("04", "Integrate ViT emotion detection (7 classes) alongside every prediction"),
        ("05", "Secure platform via JWT auth, Fernet password encryption, rate limiting & admin controls"),
        ("06", "Hybrid persistence: MongoDB 4.6 (6 collections) + SQLite; user history & batch records"),
        ("07", "Analytics dashboard with Chart.js visualisations + IsSuperUser-gated admin panel"),
        ("08", "Validate all 13 modules with a structured 44-case test catalogue (TC-AU to TC-AD)"),
    ]

    card_w  = Inches(6.15)
    card_h  = Inches(1.2)
    start_x = [Inches(0.35), Inches(6.75)]
    start_y = Inches(1.05)
    row_gap = Inches(0.1)

    for i, (num, obj) in enumerate(objectives):
        col = i % 2
        row = i // 2
        x = start_x[col]
        y = start_y + row * (card_h + row_gap)

        add_rect(slide, x, y, card_w, card_h, CARD, border_color=BORDER, border_pt=0.75)
        badge_col = PURPLE if col == 0 else GREEN
        add_rect(slide, x + Inches(0.12), y + Inches(0.28),
                 Inches(0.38), Inches(0.38), badge_col)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.28),
                    Inches(0.38), Inches(0.38),
                    num, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.6), y + Inches(0.12),
                    card_w - Inches(0.72), card_h - Inches(0.24),
                    obj, font_size=12, color=TEXT, word_wrap=True)


def build_slide_04_architecture(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "System Architecture", 4, PURPLE)

    add_image(slide, ss("fig_3_1_architecture.png"), Inches(0.35), Inches(1.0), Inches(8.6))

    tiers = [
        ("Presentation Tier", GREEN,
         "Angular 19 SPA · Bootstrap 5 · Chart.js · RxJS 7.8 · TypeScript 5.7"),
        ("API / Business Tier", PURPLE,
         "Django 5.2 · DRF · SimpleJWT · Fernet · 21 REST Endpoints · CORS"),
        ("AI / ML Tier", ORANGE,
         "MiVOLO v2 · InsightFace · SAM GAN · Fast-AgingGAN · FADING Diffusion"),
        ("Persistence Tier", GREEN,
         "MongoDB 4.6 (6 collections) · SQLite · Modal Cloud GPU · Local FS"),
    ]

    rx, ry, rw, rh = Inches(9.25), Inches(1.0), Inches(3.75), Inches(1.36)
    for title, col, body in tiers:
        add_rect(slide, rx, ry, rw, rh, CARD, border_color=col, border_pt=1.5)
        add_textbox(slide, rx + Inches(0.15), ry + Inches(0.1),
                    rw - Inches(0.3), Inches(0.38),
                    title, font_size=13, bold=True, color=col)
        add_divider(slide, rx + Inches(0.15), ry + Inches(0.5),
                    rw - Inches(0.3), col, 1)
        add_textbox(slide, rx + Inches(0.15), ry + Inches(0.56),
                    rw - Inches(0.3), Inches(0.72),
                    body, font_size=11, color=MUTED, word_wrap=True)
        ry += rh + Inches(0.08)


def build_slide_05_prediction(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Age Prediction Module", 5, PURPLE)

    lx, lw = Inches(0.35), Inches(6.0)

    add_textbox(slide, lx, Inches(1.08), lw, Inches(0.4),
                "MiVOLO v2 + InsightFace Ensemble",
                font_size=15, bold=True, color=GREEN)
    add_divider(slide, lx, Inches(1.5), lw, PURPLE, 1.5)

    add_card(slide, lx, Inches(1.58), lw, Inches(4.3),
             items=[
                 "Primary: MiVOLO v2  (ViT + YOLOv8 body context)",
                 "MAE: 3.65 yr on IMDB-WIKI / UTKFace test split",
                 "Fallback: InsightFace buffalo_l  (RetinaFace + ArcFace, MAE ~8.5)",
                 "Joint face + body bounding-box inference per image",
                 "Multi-face: detects & ages N faces simultaneously",
                 "Confidence: 1 − (σ / max_age) from ensemble spread",
                 "Emotion: trpakov/vit-face-expression  (7 classes, ViT)",
                 "Camera mode: live webcam capture + instant prediction",
             ], body_size=12.5)

    # MAE highlight badge
    add_rect(slide, lx, Inches(5.98), lw, Inches(0.95), PURPLE)
    add_textbox(slide, lx, Inches(6.08), lw, Inches(0.65),
                "MAE  3.65 yrs  |  40–50% better than single-model baselines",
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rx, rw = Inches(6.55), Inches(6.4)
    add_image(slide, ss("fig_4_3_mivolo_pipeline.png"), rx, Inches(1.08), rw)
    add_image(slide, ss("fig_4_2_age_prediction.png"),  rx, Inches(4.35), rw, Inches(2.55))
    add_textbox(slide, rx, Inches(6.98), rw, Inches(0.2),
                "MiVOLO v2 Pipeline  /  Prediction UI  (Age + Gender + Emotion overlay)",
                font_size=9, color=MUTED, align=PP_ALIGN.CENTER)


def build_slide_06_progression(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Age Progression Module", 6, ORANGE)

    add_textbox(slide, Inches(0.35), Inches(1.0), Inches(12.0), Inches(0.38),
                "4-Engine Cascade Architecture  —  Automatic quality-first fallback chain",
                font_size=14, bold=True, color=GREEN)
    add_divider(slide, Inches(0.35), Inches(1.42), Inches(12.63), ORANGE, 1.5)

    engines = [
        ("SAM GAN\n(Indian)", PURPLE, "PRIMARY",
         "pSp encoder + StyleGAN2\nIndian-finetuned checkpoint\nTarget age: 0–100\nFaceNet identity threshold 0.6"),
        ("SAM FFHQ",          PURPLE, "FALLBACK 1",
         "pSp + StyleGAN2\nGeneral-population\nWestern & diverse faces\n3–8 s GPU latency"),
        ("Fast-AgingGAN",     GREEN,  "FALLBACK 2",
         "CycleGAN architecture\n~11 MB model size\nCPU-friendly, no GPU needed\nYoung → old direction"),
        ("FADING\nDiffusion", ORANGE, "PREMIUM",
         "Stable Diffusion v1.5\nNull-text inversion\nBidirectional aging\nModal A10G GPU optional"),
    ]

    ew, eh, egap = Inches(3.05), Inches(2.1), Inches(0.23)
    for i, (name, col, badge, body) in enumerate(engines):
        x = Inches(0.35) + i * (ew + egap)
        y = Inches(1.55)
        add_rect(slide, x, y, ew, eh, CARD, border_color=col, border_pt=1.5)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.1),
                    ew - Inches(0.24), Inches(0.5),
                    name, font_size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_divider(slide, x + Inches(0.12), y + Inches(0.6),
                    ew - Inches(0.24), col, 1)
        add_textbox(slide, x + Inches(0.12), y + Inches(0.68),
                    ew - Inches(0.24), Inches(1.1),
                    body, font_size=10.5, color=MUTED, word_wrap=True)
        add_rect(slide, x + Inches(0.12), y + eh - Inches(0.34),
                 ew - Inches(0.24), Inches(0.28), col)
        add_textbox(slide, x + Inches(0.12), y + eh - Inches(0.34),
                    ew - Inches(0.24), Inches(0.28),
                    badge, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            add_rect(slide, x + ew + Inches(0.04), y + eh / 2 - Pt(1),
                     Inches(0.19), Pt(2), MUTED)

    by = Inches(3.75)
    add_image(slide, ss("fig_4_4_sam_gan_flow.png"), Inches(0.35), by, Inches(6.3))

    mx, mw, my = Inches(6.85), Inches(6.1), by + Inches(0.1)
    for col, text in [
        (GREEN,  "95%+  Identity Preservation  (FaceNet similarity ≥ 0.6)"),
        (PURPLE, "3–8 s  Progression Latency  (SAM GAN on GPU)"),
        (ORANGE, "Target Age Range: 20 · 30 · 40 · 50 · 60 · 70 · 80 yrs"),
    ]:
        add_rect(slide, mx, my, mw, Inches(0.88), col)
        add_textbox(slide, mx, my, mw, Inches(0.88),
                    text, font_size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        my += Inches(1.0)

    add_textbox(slide, Inches(0.35), Inches(7.16), Inches(6.3), Inches(0.2),
                "SAM GAN pSp Encoder → StyleGAN2 W+ Decoder Pipeline",
                font_size=9, color=MUTED, align=PP_ALIGN.CENTER)


def build_slide_07_demo(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Application Demo — Angular 19 SPA", 7, GREEN)

    add_textbox(slide, Inches(0.35), Inches(0.98), Inches(12.5), Inches(0.3),
                "21 REST Endpoints  •  Dark / Light Themes  •  Mobile-Responsive  •  JWT-Secured",
                font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

    screenshots = [
        (ss("fig_4_8_dashboard.png"),            Inches(0.35), Inches(1.35), Inches(6.1),
         "Dashboard — KPI Cards & Quick Actions", GREEN),
        (ss("fig_4_2_age_prediction.png"),        Inches(6.75), Inches(1.35), Inches(6.2),
         "Age Prediction — Upload / Camera Modes", PURPLE),
        (ss("fig_4_12_analysis_dashboard_1.png"), Inches(0.35), Inches(4.45), Inches(6.1),
         "Analytics — Charts, Distributions", ORANGE),
        (ss("fig_4_14_admin_panel.png"),           Inches(6.75), Inches(4.45), Inches(6.2),
         "Admin Panel — User Management & Health", GREEN),
    ]

    for path, x, y, w, caption, col in screenshots:
        add_image(slide, path, x, y, w, Inches(2.9))
        # coloured border only — no fill overlay
        add_rect(slide, x, y, w, Inches(2.9), RGBColor(0xFF, 0xFF, 0xFF),
                 border_color=col, border_pt=1.5)
        add_textbox(slide, x, y + Inches(2.95), w, Inches(0.28),
                    caption, font_size=10, bold=True, color=col, align=PP_ALIGN.CENTER)


def build_slide_08_results(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Results & Performance", 8, PURPLE)

    lx, lw = Inches(0.35), Inches(5.9)

    add_textbox(slide, lx, Inches(1.08), lw, Inches(0.35),
                "Age Prediction Accuracy", font_size=13, bold=True, color=PURPLE)

    tbl1 = slide.shapes.add_table(4, 3, lx, Inches(1.48), lw, Inches(2.0)).table
    for r, row in enumerate([
        ["Model", "MAE (years)", "Use Case"],
        ["MiVOLO v2 (Primary)",      "3.65", "Main engine"],
        ["InsightFace buffalo_l",     "~8.5", "Fallback"],
        ["DeepFace EfficientNet (ref)", "~6.0", "Baseline"],
    ]):
        for c, val in enumerate(row):
            _add_table_cell(tbl1, r, c, val)
    style_table(tbl1, PURPLE,
                RGBColor(0xED, 0xE8, 0xFC), RGBColor(0xF8, 0xF6, 0xFF),
                body_text=TEXT, hdr_pt=12, body_pt=11)

    add_textbox(slide, lx, Inches(3.6), lw, Inches(0.35),
                "Age Progression Quality", font_size=13, bold=True, color=GREEN)

    tbl2 = slide.shapes.add_table(5, 3, lx, Inches(4.0), lw, Inches(2.25)).table
    for r, row in enumerate([
        ["Engine",           "Identity Score", "Latency"],
        ["SAM GAN (Indian)", "0.82 avg",       "4–6 s"],
        ["SAM FFHQ",         "0.79 avg",       "4–6 s"],
        ["Fast-AgingGAN",    "0.71 avg",       "1–2 s"],
        ["FADING Diffusion", "0.88 avg",       "6–10 s GPU"],
    ]):
        for c, val in enumerate(row):
            _add_table_cell(tbl2, r, c, val)
    style_table(tbl2, GREEN,
                RGBColor(0xE8, 0xF7, 0xF0), RGBColor(0xF4, 0xFC, 0xF8),
                body_text=TEXT, hdr_pt=12, body_pt=11)

    add_image(slide, ss("Figure_4_6_Comparison.png"),
              Inches(6.45), Inches(1.08), Inches(6.6))
    add_textbox(slide, Inches(6.45), Inches(7.06), Inches(6.6), Inches(0.2),
                "Progression Comparison Grid — Multiple ethnicities, ages 30 → 70",
                font_size=9, color=MUTED, align=PP_ALIGN.CENTER)

    add_rect(slide, lx, Inches(6.4), lw, Inches(0.62), PURPLE)
    add_textbox(slide, lx, Inches(6.5), lw, Inches(0.42),
                "Best result: Korean Male @ Age 40 →  97.0% Accuracy  |  SSIM 0.92  |  PSNR 27.74 dB",
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_09_testing(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Testing & Validation", 9, GREEN)

    tbl = slide.shapes.add_table(
        9, 4, Inches(0.35), Inches(1.08), Inches(12.63), Inches(4.9)
    ).table

    for r, row in enumerate([
        ["Test Suite",        "Coverage Area",                          "Cases", "Status"],
        ["TC-AUTH-001–006",   "Authentication & JWT Token Management",  "6",     "PASS"],
        ["TC-PRED-001–010",   "Single & Multi-Face Age Prediction",     "10",    "PASS"],
        ["TC-BATCH-001–005",  "Batch Upload & Processing Pipeline",     "5",     "PASS"],
        ["TC-PROG-001–008",   "GAN / Diffusion Age Progression",        "8",     "PASS"],
        ["TC-HIST-001–004",   "Prediction History & MongoDB Retrieval", "4",     "PASS"],
        ["TC-ANAL-001–006",   "Analytics Dashboard & Chart Data",       "6",     "PASS"],
        ["TC-ADMIN-001–005",  "Admin Panel: Users, Health, Logs",       "5",     "PASS"],
        ["TOTAL",             "All 13 Functional Modules",              "44",    "100% PASS"],
    ]):
        for c, val in enumerate(row):
            _add_table_cell(tbl, r, c, val)
    style_table(tbl, PURPLE,
                RGBColor(0xED, 0xE8, 0xFC), RGBColor(0xF8, 0xF6, 0xFF),
                body_text=TEXT, hdr_pt=13, body_pt=12,
                last_row_fill=GREEN, last_row_text=WHITE)

    for col, bx, text in [
        (PURPLE, Inches(0.35),  "44 Test Cases  |  100% Pass Rate"),
        (GREEN,  Inches(4.6),   "13 Modules Fully Implemented"),
        (ORANGE, Inches(8.85),  "21 REST API Endpoints Validated"),
    ]:
        add_rect(slide, bx, Inches(6.1), Inches(3.93), Inches(0.95), col)
        add_textbox(slide, bx, Inches(6.25), Inches(3.93), Inches(0.65),
                    text, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_10_conclusion(prs, layout):
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide)
    add_slide_header(slide, "Conclusion & Future Work", 10, ORANGE)

    add_card(slide, Inches(0.35), Inches(1.08), Inches(6.1), Inches(6.1),
             title="Key Achievements", title_color=GREEN,
             items=[
                 "MiVOLO v2 ensemble — MAE 3.65 yrs (40-50% better than baselines)",
                 "4-engine progression cascade with 95%+ identity preservation",
                 "FADING diffusion — bidirectional photorealistic aging",
                 "13 functional modules, 21 REST endpoints, JWT-secured",
                 "Hybrid MongoDB + SQLite persistence (6 collections)",
                 "Full analytics dashboard + IsSuperUser admin panel",
                 "Glassmorphic Angular 19 SPA — dark & light themes",
             ], body_size=13)

    add_rect(slide, Inches(6.65), Inches(1.08), Inches(6.33), Inches(6.1),
             CARD, border_color=ORANGE, border_pt=1.5)
    add_textbox(slide, Inches(6.82), Inches(1.18), Inches(6.0), Inches(0.38),
                "Future Work", font_size=13, bold=True, color=ORANGE)
    add_divider(slide, Inches(6.82), Inches(1.58), Inches(6.0), PURPLE, 1.5)

    fy = Inches(1.72)
    for num, text in [
        ("01", "3D face mesh aging (Neural Radiance Fields / NeRF)"),
        ("02", "Mobile-optimised TFLite / CoreML export for iOS & Android"),
        ("03", "Real-time video stream age-overlay pipeline"),
        ("04", "Federated learning for privacy-preserving model retraining"),
        ("05", "Expanded Indian ethnic dataset for SAM GAN fine-tuning"),
        ("06", "Cross-platform Electron desktop & PDF batch-export module"),
    ]:
        add_rect(slide, Inches(6.82), fy, Inches(0.42), Inches(0.38), ORANGE)
        add_textbox(slide, Inches(6.82), fy, Inches(0.42), Inches(0.38),
                    num, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(7.34), fy, Inches(5.48), Inches(0.38),
                    text, font_size=12.5, color=TEXT, word_wrap=True)
        fy += Inches(0.84)

    add_textbox(slide, Inches(0.35), Inches(7.11), Inches(12.63), Inches(0.24),
                "Thank You  —  Open to Questions",
                font_size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════ #
#                                   MAIN                                      #
# ═══════════════════════════════════════════════════════════════════════════ #

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    print("Building slides...")
    build_slide_01_title(prs, blank);      print("  [1/10] Title")
    build_slide_02_problem(prs, blank);    print("  [2/10] Problem Statement")
    build_slide_03_objectives(prs, blank); print("  [3/10] Objectives")
    build_slide_04_architecture(prs, blank); print("  [4/10] Architecture")
    build_slide_05_prediction(prs, blank); print("  [5/10] Age Prediction")
    build_slide_06_progression(prs, blank); print("  [6/10] Age Progression")
    build_slide_07_demo(prs, blank);       print("  [7/10] App Demo")
    build_slide_08_results(prs, blank);    print("  [8/10] Results")
    build_slide_09_testing(prs, blank);    print("  [9/10] Testing")
    build_slide_10_conclusion(prs, blank); print("  [10/10] Conclusion")

    out = OUT
    try:
        prs.save(out)
    except PermissionError:
        out = OUT.replace(".pptx", "_new.pptx")
        prs.save(out)
        print(f"WARNING: Original file was open. Saved to: {out}")

    print(f"\nSaved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
